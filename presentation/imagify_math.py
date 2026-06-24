# -*- coding: utf-8 -*-
"""
For the math-heavy slides that the PDF->PPTX converter scatters, replace the
slide BODY with the crisp polished-LaTeX render (formulas 100% correct, zero
scatter), keeping a NATIVE editable teal title bar on top. Other slides stay
fully editable.
"""
import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

PPTX = "SE2RoadNet_Test_Prioritization_SDC_editable.pptx"
BODYDIR = "qa/bodyimg"
TEAL = RGBColor(0x23, 0x37, 0x3B)
WHITE = RGBColor(0xFA, 0xFA, 0xFA)

# slide index (1-based) -> (clean title, polished-pdf page render path)
TARGETS = {
    5:  "Phát biểu bài toán & Tiêu chí đánh giá",
    18: "Ý tưởng cốt lõi: Bất biến SE(2)",
    19: "Bước 1: 7 kênh đặc trưng bất biến SE(2)",
    21: "Bước 3: Attention với bias hiệu số ∆s",
    22: "Bước 4: Huấn luyện",
    24: "Chỉ số đánh giá: APFD",
    33: "Backup: APFD — ví dụ tính tay",
}

def crop_body(src, dst):
    """Crop off the top dark title-bar strip; return body aspect (w/h)."""
    im = Image.open(src).convert("RGB")
    w, h = im.size
    px = im.load()
    # find first mostly-light row after the dark top bar
    def row_dark(y):
        dark = 0
        for x in range(0, w, 40):
            r, g, b = px[x, y]
            if r < 80 and g < 90 and b < 90:
                dark += 1
        return dark > (w // 40) * 0.6
    y = 0
    # bar must start dark near top
    if row_dark(2):
        while y < h * 0.25 and row_dark(y):
            y += 1
        y += 2  # tiny pad
    else:
        y = int(h * 0.118)
    body = im.crop((0, y, w, h))
    body.save(dst)
    return body.size, y / h

os.makedirs(BODYDIR, exist_ok=True)
prs = Presentation(PPTX)
SW, SH = prs.slide_width, prs.slide_height
sw_in, sh_in = Emu(SW).inches, Emu(SH).inches

for idx, title in TARGETS.items():
    src = f"{BODYDIR}/p-{idx:02d}.png"
    dst = f"{BODYDIR}/body-{idx:02d}.png"
    (bw, bh), barfrac = crop_body(src, dst)
    bar_h = barfrac * sh_in            # native bar height matches the crop
    slide = prs.slides[idx - 1]
    # delete all shapes
    for sh in list(slide.shapes):
        sh._element.getparent().remove(sh._element)
    # teal title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(bar_h))
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()
    bar.shadow.inherit = False
    tf = bar.text_frame; tf.word_wrap = False
    tf.margin_left = Inches(0.14); tf.margin_top = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    r.font.bold = True; r.font.size = Pt(15); r.font.color.rgb = WHITE
    # body image: fill area below bar, fit by height, center horizontally
    avail_w, avail_h = sw_in, sh_in - bar_h
    asp = bw / bh
    img_h = avail_h
    img_w = img_h * asp
    if img_w > avail_w:
        img_w = avail_w; img_h = img_w / asp
    left = (sw_in - img_w) / 2
    top = bar_h + (avail_h - img_h) / 2
    slide.shapes.add_picture(dst, Inches(left), Inches(top),
                             Inches(img_w), Inches(img_h))
    print(f"slide {idx}: bar={bar_h:.3f}in body={bw}x{bh} -> {img_w:.2f}x{img_h:.2f}in")

prs.save(PPTX)
print("saved", PPTX)
