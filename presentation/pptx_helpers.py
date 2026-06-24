# -*- coding: utf-8 -*-
"""Native python-pptx helpers for the SE2RoadNet deck (16:9, 13.333x7.5in)."""
import math
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from PIL import Image

W, H = 13.333, 7.5
TEAL   = RGBColor(0x23, 0x37, 0x3B)
ORANGE = RGBColor(0xEB, 0x81, 0x1B)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)
RED    = RGBColor(0xB2, 0x3A, 0x33)
ALERT  = RGBColor(0xB0, 0x30, 0x30)
GREY   = RGBColor(0x8A, 0x97, 0x9B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
NAVYMID= RGBColor(0x5B, 0x7A, 0x82)
SOFT   = RGBColor(0xF2, 0xF4, 0xF7)
TINT   = RGBColor(0xD6, 0xDE, 0xE0)
FONT   = "Calibri"
_AL = {'l': PP_ALIGN.LEFT, 'c': PP_ALIGN.CENTER, 'r': PP_ALIGN.RIGHT}
_AN = {'t': MSO_ANCHOR.TOP, 'm': MSO_ANCHOR.MIDDLE, 'b': MSO_ANCHOR.BOTTOM}


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _runs(p, content, size, color, bold=False, italic=False):
    """content: str or list of run dicts {t,b,i,c,sz}."""
    items = [{'t': content}] if isinstance(content, str) else content
    for rd in items:
        r = p.add_run(); r.text = rd.get('t', '')
        r.font.size = Pt(rd.get('sz', size))
        r.font.bold = rd.get('b', bold)
        r.font.italic = rd.get('i', italic)
        r.font.color.rgb = rd.get('c', color)
        r.font.name = FONT


def text(slide, x, y, w, h, lines, size=15, color=TEAL, align='l', anchor='t',
         bold=False, italic=False, sp_after=4, line_sp=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = _AN[anchor]
    for m in ('left', 'right'): setattr(tf, f'margin_{m}', Pt(2))
    for m in ('top', 'bottom'): setattr(tf, f'margin_{m}', Pt(1))
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.alignment = _AL[align]; p.space_after = Pt(sp_after)
        if line_sp: p.line_spacing = line_sp
        _runs(p, ln, size, color, bold, italic)
    return tb


def bar(slide, title, page=None, size=26):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(W), Inches(0.92))
    r.fill.solid(); r.fill.fore_color.rgb = TEAL
    r.line.fill.background(); r.shadow.inherit = False
    tf = r.text_frame; tf.word_wrap = False; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.4)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run(); run.text = title
    run.font.bold = True; run.font.size = Pt(size); run.font.color.rgb = WHITE; run.font.name = FONT
    if page is not None: footer(slide, page)
    return r


def footer(slide, txt):
    text(slide, W - 1.3, H - 0.48, 1.1, 0.35, [str(txt)], size=11, color=GREY, align='r')


def bullets(slide, x, y, w, h, items, size=15, color=TEAL, gap=7):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        txt, lvl = it if isinstance(it, tuple) else (it, 0)
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.level = lvl; p.space_after = Pt(gap)
        rm = p.add_run(); rm.text = "▸  " if lvl == 0 else "–  "
        rm.font.color.rgb = ORANGE if lvl == 0 else NAVYMID
        rm.font.size = Pt(size); rm.font.bold = True; rm.font.name = FONT
        _runs(p, txt, size, color)
    return tb


def box(slide, x, y, w, h, fill=None, line=None, rounded=True, width=1.2):
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = line; s.line.width = Pt(width)
    s.shadow.inherit = False
    s.text_frame.clear()
    return s


def card(slide, x, y, w, h, title, body, fill, line, tcolor=None, size=12, title_size=None):
    box(slide, x, y, w, h, fill, line)
    lines = [[{'t': title, 'b': True, 'c': tcolor or line, 'sz': title_size or size + 3}]]
    for ln in (body if isinstance(body, list) else [body]):
        lines.append(ln)
    text(slide, x + 0.14, y + 0.10, w - 0.28, h - 0.20, lines, size=size, color=TEAL, anchor='t')


def callout(slide, x, y, w, h, content, fill=None, line=ORANGE, color=TEAL, size=14, align='c'):
    """orange takeaway-style box."""
    box(slide, x, y, w, h, fill if fill is not None else RGBColor(0xFD, 0xF1, 0xE3), line)
    text(slide, x + 0.15, y, w - 0.30, h, content if isinstance(content, list) else [content],
         size=size, color=color, align=align, anchor='m', bold=True)


def table(slide, x, y, w, rows, col_w, fontsize=12, header=True, hi=None,
          row_h=0.30, aligns=None):
    nr, nc = len(rows), len(rows[0])
    gt = slide.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w),
                                Inches(row_h * nr)).table
    gt.first_row = False; gt.horz_banding = False
    for ci, cw in enumerate(col_w):
        gt.columns[ci].width = Inches(cw)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            c = gt.cell(ri, ci)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = Pt(5); c.margin_right = Pt(3)
            c.margin_top = Pt(1); c.margin_bottom = Pt(1)
            tf = c.text_frame; p = tf.paragraphs[0]
            p.alignment = _AL[(aligns[ci] if aligns else 'l')]
            r = p.add_run(); r.text = str(val); r.font.size = Pt(fontsize); r.font.name = FONT
            ishdr = (ri == 0 and header)
            r.font.bold = ishdr or (ri == hi)
            r.font.color.rgb = WHITE if ishdr else TEAL
            c.fill.solid()
            c.fill.fore_color.rgb = TEAL if ishdr else (RGBColor(0xE7, 0xF1, 0xE7) if ri == hi else WHITE)
    return gt


def arrow(slide, x1, y1, x2, y2, color=TEAL, width=2.0, dashed=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(width); c.shadow.inherit = False
    ln = c.line._get_or_add_ln()
    if dashed:
        d = ln.makeelement(qn('a:prstDash'), {'val': 'dash'}); ln.append(d)
    end = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(end)
    return c


def rule(slide, x, y, w, color, width=1.0):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y), Inches(x + w), Inches(y))
    c.line.color.rgb = color; c.line.width = Pt(width); c.shadow.inherit = False
    return c


def badge(slide, cx, cy, d, num, fill, tcolor=WHITE, size=12):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d/2), Inches(cy - d/2), Inches(d), Inches(d))
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.fill.background(); s.shadow.inherit = False
    tf = s.text_frame; tf.word_wrap = False; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(num); r.font.bold = True; r.font.color.rgb = tcolor; r.font.size = Pt(size); r.font.name = FONT
    return s


def chip(slide, x, y, w, h, content, line=TEAL, fill=WHITE, color=TEAL, size=13, bold=True):
    """small bordered label box (fcolorbox-style)."""
    box(slide, x, y, w, h, fill, line, rounded=False, width=1.0)
    text(slide, x, y, w, h, content if isinstance(content, list) else [content],
         size=size, color=color, align='c', anchor='m', bold=bold)


def sine(slide, x, y, w, h, color=ORANGE, width=2.6, cycles=1.3, amp=0.62, n=64,
         rot=0.0, slope=0.18):
    """native freeform sine-ish road curve inside box (x,y,w,h)."""
    cx, cy = x + w/2, y + h/2
    ct, st = math.cos(math.radians(rot)), math.sin(math.radians(rot))
    pts = []
    for i in range(n + 1):
        t = i/n
        lx = (t - 0.5) * w
        ly = -amp * (h/2) * math.sin(2*math.pi*cycles*t) - slope*lx
        rx = lx*ct - ly*st; ry = lx*st + ly*ct
        pts.append((Inches(cx + rx).emu, Inches(cy + ry).emu))
    fb = slide.shapes.build_freeform(pts[0][0], pts[0][1], scale=1.0)
    fb.add_line_segments(pts[1:], close=False)
    shp = fb.convert_to_shape()
    shp.fill.background(); shp.line.color.rgb = color; shp.line.width = Pt(width)
    shp.shadow.inherit = False
    return shp


def axis(slide, x, y, w, h, color=NAVYMID):
    arrow(slide, x, y + h, x + w, y + h, color, 1.0)   # x-axis
    arrow(slide, x, y + h, x, y, color, 1.0)           # y-axis


def imgfit(slide, path, x, y, w, h, align='c', valign='m'):
    iw, ih = Image.open(path).size; asp = iw / ih
    bw, bh = w, h
    if bw / bh > asp: bw = bh * asp
    else: bh = bw / asp
    px = x + (w - bw) * {'l': 0, 'c': 0.5, 'r': 1}[align]
    py = y + (h - bh) * {'t': 0, 'm': 0.5, 'b': 1}[valign]
    return slide.shapes.add_picture(path, Inches(px), Inches(py), Inches(bw), Inches(bh))
