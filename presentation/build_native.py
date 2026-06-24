# -*- coding: utf-8 -*-
"""Build SE2RoadNet deck NATIVELY in python-pptx (editable text + code-rendered
figures). Figures: assets/fig-*.png (deck TikZ + formulas), figures/gen_*.png
(matplotlib), pipelines, logos."""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx_helpers import (blank, bar, footer, text, bullets, box, card, callout,
                          table, imgfit, arrow, rule, badge, chip, sine, axis,
                          W, H, TEAL, ORANGE, GREEN, RED, ALERT,
                          GREY, WHITE, NAVYMID, SOFT, TINT)
ROT = RGBColor(0xB8, 0x54, 0x50)   # terracotta (blind-spot 1)
RES = RGBColor(0xC8, 0xA9, 0x51)   # gold (blind-spot 2)
GTINT = RGBColor(0xEA, 0xF1, 0xEA)
OTINT = RGBColor(0xFD, 0xF1, 0xE3)
NTINT = RGBColor(0xEE, 0xF0, 0xF5)

A = "assets/"; F = "figures/"
prs = Presentation()
prs.slide_width = Inches(W); prs.slide_height = Inches(H)

def divider(title):
    s = blank(prs)
    box(s, 0, 0, W, H, fill=WHITE, line=None, rounded=False)
    text(s, 1.2, 3.0, W-2.4, 1.0, [title], size=34, color=TEAL, align='l', bold=True, anchor='m')
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.25), Inches(4.05), Inches(2.6), Pt(3))
    ln.fill.solid(); ln.fill.fore_color.rgb = ORANGE; ln.line.fill.background(); ln.shadow.inherit=False
    return s

I = lambda c: {'t': c}
B = lambda c, col=TEAL: {'t': c, 'b': True, 'c': col}

# ============ 1. TITLE ============
s = blank(prs)
imgfit(s, F+"hcmus.png", 0.7, 0.7, 1.9, 1.9)
imgfit(s, F+"icst_logo.jpg", W-2.5, 0.9, 1.7, 1.5)
text(s, 3.0, 0.7, 7.3, 2.0, [
    [B("The Self-Driving Car Testing Competition", TEAL)],
    [I("organized jointly with the International Conference on")],
    [I("Software Testing, Verification and Validation "), B("(ICST)")],
    [{'t':"Test Prioritization for Self-Driving Cars", 'b':True, 'c':TEAL, 'sz':24}],
    [B("SE2RoadNet", ORANGE), {'t':" — bất biến SE(2) bằng kiến trúc",'i':True}],
], size=15, align='c', anchor='t', sp_after=6)
card(s, 1.4, 4.2, 4.7, 2.0, "Tóm tắt đề tài",
     [[I("Xếp hạng kịch bản kiểm thử SDC để "), B("lộ lỗi sớm",ORANGE), I("; mô hình")],
      [B("bất biến phép xoay / lấy mẫu",TEAL), I(" bằng ràng buộc kiến trúc,")],
      [I("không bằng augmentation.")]], SOFT, TINT, tcolor=TEAL, size=14)
card(s, 6.6, 4.2, 4.5, 2.0, "Thành viên nhóm",
     [[I("Trần Chí Nguyên — 23102244")],
      [I("Huỳnh Trung Kiệt — 23122039")],
      [I("Đào Sỹ Duy Minh — 23122041")]], SOFT, TINT, tcolor=TEAL, size=14)
text(s, 0, 6.6, W, 0.4, ["University of Science, VNU-HCM   ·   ICST 2026"], size=13.5, color=GREY, align='c')

# ============ 2. OUTLINE ============
s = blank(prs); bar(s, "Nội dung chính", "1/37")
outline = ["Bối cảnh & Phát biểu bài toán",
           "Công trình liên quan & Khoảng trống nghiên cứu",
           "RoadFury: Phương pháp nền tảng",
           "SE2RoadNet: Đề xuất của nhóm",
           "Đánh giá thực nghiệm",
           "Hạn chế & Hướng phát triển"]
tb = s.shapes.add_textbox(Inches(1.4), Inches(1.6), Inches(10.5), Inches(4.3)); tf = tb.text_frame; tf.word_wrap=True
for i, it in enumerate(outline):
    p = tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after = Pt(14)
    rn = p.add_run(); rn.text = f"{i+1}.  "; rn.font.bold=True; rn.font.size=Pt(20); rn.font.color.rgb=ORANGE; rn.font.name="Calibri"
    r = p.add_run(); r.text = it; r.font.size=Pt(20); r.font.name="Calibri"
    r.font.bold = (i==3); r.font.color.rgb = ORANGE if i==3 else TEAL
text(s, 1.4, 6.3, 10.5, 0.5,
     [[{'t':"(+ Backup slides cho Q&A: APFD worked-example, Multi-trial, AUC vs APFD, Resolution probe, Chi phí ở quy mô lớn)",'i':True,'c':GREY,'sz':12}]], align='l')

# ============ 3. divider ============
divider("Bối cảnh & Phát biểu bài toán")

# ============ 4. Bối cảnh (NATIVE) ============
s = blank(prs); bar(s, "Bối cảnh: Kiểm thử Self-Driving Car trong mô phỏng", "3/37")
ctx = [(0.5, "Hàng ngàn kịch bản", TEAL, NTINT, "Sinh tự động các con đường (ambiegen, frenetic…) trên BeamNG.tech. Mỗi đường = một bài kiểm thử giữ làn."),
       (4.65, "Mô phỏng tốn kém", ORANGE, OTINT, "Mỗi kịch bản chạy vật lý thời gian thực ⇒ nhiều giờ CPU cho toàn tập."),
       (8.8, "Phải ưu tiên", GREEN, GTINT, "Xếp thứ tự để lỗi lộ sớm ⇒ chạy ít kịch bản đầu là đủ, cắt 50–80% thời gian.")]
for x, ttl, c, fill, body in ctx:
    box(s, x, 1.2, 3.85, 1.95, fill, c)
    text(s, x+0.18, 1.32, 3.5, 1.7, [[B(ttl, c)], [I(body)]], size=13.5, sp_after=4)
arrow(s, 4.42, 2.15, 4.62, 2.15, NAVYMID, 2.4)
arrow(s, 8.57, 2.15, 8.77, 2.15, NAVYMID, 2.4)
sine(s, 1.0, 3.5, 11.3, 1.6, color=RGBColor(0xC9,0xD5,0xD8), width=11, cycles=1.7, amp=0.5, slope=0)
sine(s, 1.0, 3.5, 11.3, 1.6, color=WHITE, width=1.4, cycles=1.7, amp=0.5, slope=0)
text(s, 2.0, 3.3, 2.0, 0.3, [[B("✓ PASS", GREEN)]], size=13.5)
text(s, 8.5, 4.55, 2.0, 0.3, [[B("✗ FAIL", ALERT)]], size=13.5)
d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.6), Inches(4.25), Inches(0.14), Inches(0.14))
d.fill.solid(); d.fill.fore_color.rgb=ORANGE; d.line.fill.background(); d.shadow.inherit=False
callout(s, 1.2, 6.2, 10.9, 0.85,
        ["Một thay đổi nhỏ có thể khiến xe lệch làn ở kịch bản hiếm ⇒ cần kiểm thử quy mô lớn, nhưng chạy hết thì quá đắt."],
        size=15)

# ============ 5. Phát biểu bài toán (NATIVE) ============
s = blank(prs); bar(s, "Phát biểu bài toán & Tiêu chí đánh giá", "4/37")
box(s, 0.5, 1.15, 3.5, 2.0, RGBColor(0xEE,0xF0,0xF5), TEAL)
text(s, 0.65, 1.25, 3.2, 1.8, [
    [B("INPUT", TEAL)],
    [I("• Đường = chuỗi điểm 2D")],
    [I("  R = {p₁,…,pₙ} ⊂ ℝ²")],
    [I("• N thay đổi: 64–197 điểm")],
    [I("• Chỉ hình học, không hành vi xe")]], size=13, sp_after=2)
box(s, 5.0, 1.15, 3.0, 2.0, OTINT, ORANGE)
text(s, 5.0, 1.15, 3.0, 2.0, [[B("SCORER fθ", ORANGE)], [I("R ↦ [0,1]")], [I("(xác suất FAIL)")]], size=14, align='c', anchor='m', sp_after=4)
box(s, 9.0, 1.15, 3.8, 2.0, GTINT, GREEN)
text(s, 9.15, 1.25, 3.5, 1.8, [
    [B("OUTPUT", GREEN)],
    [I("• Hoán vị π của n kịch bản")],
    [I("• Sắp giảm dần theo điểm")],
    [I("• Mục tiêu:")],
    [I("  π* = argmax APFD(π)")]], size=13, sp_after=2)
arrow(s, 4.05, 2.15, 4.95, 2.15, TEAL, 2.6)
arrow(s, 8.05, 2.15, 8.95, 2.15, TEAL, 2.6)
box(s, 1.8, 3.55, 9.7, 0.7, GTINT, GREEN)
text(s, 1.8, 3.55, 9.7, 0.7, [[B("Success Criterion: ",GREEN),I("lỗi (FAIL) được xếp càng sớm càng tốt ⇒ APFD cao ⇒ tiết kiệm chi phí mô phỏng")]], size=13.5, align='c', anchor='m')
cols = [(0.6, "Cho trước", TEAL, ["Nhãn train y ∈ {0,1} (PASS/FAIL)", "Cùng ngân sách cho mọi tool"]),
        (4.9, "Khó khăn", ALERT, ["Chỉ có hình học đường, không có phản ứng xe", "Tập thi out-of-distribution"]),
        (9.2, "Metric", ORANGE, ["APFD (primary)", "AUC (phụ), Rotation Δ"])]
for x, ttl, c, items in cols:
    text(s, x, 5.1, 3.9, 0.4, [[B(ttl, c)]], size=15)
    bullets(s, x, 5.5, 3.9, 1.7, items, size=13.5, gap=4)

# ============ 6. Vì sao ưu tiên ============
s = blank(prs); bar(s, "Vì sao cần ưu tiên kiểm thử?", "5/37")
imgfit(s, F+"gen_apfd_curve.png", 0.5, 1.2, 7.0, 4.6, valign='t')
text(s, 7.8, 1.3, 5.1, 0.4, [[B("Trực giác:")]], size=16)
bullets(s, 7.8, 1.8, 5.1, 3.2, [
    [I("Xếp hạng "), B("tốt",GREEN), I(" ⇒ lỗi dồn lên đầu ⇒ đường cong dốc đứng từ sớm.")],
    [I("Xếp hạng "), B("kém",ALERT), I(" ⇒ lỗi nằm cuối hàng đợi ⇒ tốn nhiều giờ mô phỏng.")],
    [B("APFD"), I(" đo chính diện tích dưới đường cong này.")]], size=15, gap=9)
callout(s, 7.8, 4.9, 5.1, 1.0, [[I("956 kịch bản × mô phỏng vật lý")],[B("⇒ nhiều giờ CPU nếu chạy hết", ALERT)]],
        line=ALERT, size=14)
callout(s, 1.0, 6.35, 11.3, 0.8,
        ["Ranking tốt cắt 50–80% chi phí testing — bài toán không phải “chạy nhanh hơn” mà là “chạy đúng thứ tự”."], size=15)

# ============ 7. Dataset ============
s = blank(prs); bar(s, "Dataset: SensoDat (Birchler et al., MSR 2024)", "6/37")
text(s, 0.5, 1.15, 6.2, 0.4, [[B("Tổng quan — 3 split trên BeamNG.tech:")]], size=15)
bullets(s, 0.5, 1.6, 6.2, 1.9, [
    [B("Train"), I(" – huấn luyện mô hình.")],
    [B("Test (SensoDat)"), I(" – đánh giá in-distribution.")],
    [B("Competition"), I(" – 956 kịch bản "), B("out-of-distribution",ALERT), I(".")]], size=14, gap=5)
table(s, 0.5, 3.6, 6.2, [
    ["Split","N","%FAIL","Dài (m)"],
    ["Train","28 804","38.4%","61–454"],
    ["Test (SensoDat)","7 202","38.4%","76–433"],
    ["Competition (OOD)","956","36.9%","129–229"]],
    [2.6,1.3,1.2,1.1], fontsize=13.5, hi=3, aligns=['l','r','r','r'])
imgfit(s, F+"gen_failrate_by_length.png", 7.0, 1.2, 5.9, 4.4, valign='t')
text(s, 7.0, 5.7, 5.9, 0.9, [[{'t':"Competition có khoảng độ dài hẹp hơn (129–229m) ⇒ dịch phân phối (distribution shift).",'i':True,'c':GREY}]], size=13.5, align='c')

# ============ 8. Ví dụ ============
s = blank(prs); bar(s, "Ví dụ kịch bản: hình dạng quyết định nhãn", "7/37")
imgfit(s, F+"real_roads_grid.png", 0.4, 1.15, 6.4, 3.7, valign='t')
text(s, 0.4, 4.9, 6.4, 0.4, [[{'t':"Đường thật từ SensoDat test split  ·  ",'i':True,'c':GREY},B("PASS",GREEN),I(" mượt, "),B("FAIL",ALERT),I(" chicane gắt")]], size=13.5, align='c')
imgfit(s, F+"example.png", 7.0, 1.15, 5.9, 2.1, valign='t')
text(s, 7.0, 3.25, 5.9, 0.4, [[{'t':"BeamNG.tech: trái PASS (trong làn) · phải FAIL (lệch làn)",'i':True,'c':GREY}]], size=12.5, align='c')
imgfit(s, F+"gen_curvature_signature.png", 7.4, 3.65, 5.1, 2.0, valign='t')
callout(s, 0.9, 6.35, 11.5, 0.8, [[I("⇒ Không có toạ độ (x, y) tuyệt đối nào quyết định nhãn — chỉ "),B("hình dạng",TEAL),I(". Gợi ý cho ràng buộc "),B("bất biến",ORANGE),I(".")]], size=15)

# ============ 9. divider ============
divider("Công trình liên quan & Khoảng trống nghiên cứu")

# ============ 10. RW 3 cards (NATIVE) ============
s = blank(prs); bar(s, "Related Work: 3 hướng tiếp cận", "8/37")
def rwcard(x, num, title, tcol, fill, line, idea, pro, con):
    box(s, x, 1.3, 3.85, 4.55, fill, line)
    badge(s, x+0.42, 1.72, 0.42, num, line)
    text(s, x+0.7, 1.5, 3.0, 0.7, [[B(title, tcol)]], size=14)
    rule(s, x+0.2, 2.25, 3.45, line)
    text(s, x+0.2, 2.4, 3.45, 1.9, [[B("Ý tưởng: "),I(idea)]], size=13)
    text(s, x+0.2, 4.35, 3.45, 0.7, [[B("✓ Ưu: ",GREEN),I(pro)]], size=13)
    text(s, x+0.2, 5.05, 3.45, 0.75, [[B("✗ Hạn chế: ",ALERT),I(con)]], size=13)
rwcard(0.45, 1, "Search-based / Diversity", TEAL, SOFT, TEAL,
       "GA / heuristic chọn tập test đa dạng hành vi (SO-SDC-Prioritizer, Greedy-diversity).",
       "không cần nhãn, truyền thống SBST.",
       "đa dạng ≠ lộ lỗi; đặc trưng phụ thuộc khung.")
rwcard(4.62, 2, "Feature-based ML", ORANGE, OTINT, ORANGE,
       "vài đặc trưng tay → classifier nhỏ → rank (ITEP4SDC: MLP 3 đặc trưng).",
       "rẻ, dễ giải thích.",
       "ít đặc trưng, phụ thuộc khung, tinh chỉnh theo từng bench.")
rwcard(8.79, 3, "Deep: GNN / CNN / Transformer", GREEN, GTINT, GREEN,
       "học biểu diễn từ đồ thị đường / ảnh render / chuỗi điểm (GNN, ResNet, RoadFury).",
       "APFD cao nhất (RoadFury 0.804).",
       "không bất biến — xoay / lấy mẫu là điểm trôi.")
callout(s, 0.7, 6.15, 11.9, 0.85,
        [[B("⇒ Khoảng trống: ",TEAL),I("chưa phương pháp nào đảm bảo bất biến với phép xoay / tần số lấy mẫu — tất cả dựa vào augmentation (xấp xỉ) hoặc đặc trưng phụ thuộc khung.")]], size=14)

# ============ 11. RW table+scatter ============
s = blank(prs); bar(s, "Related Work: Tổng hợp & Khoảng trống", "9/37")
table(s, 0.5, 1.2, 6.3, [
    ["Method","Approach","APFD","Bất biến"],
    ["Random","–","0.493","✗"],
    ["GNN","graph","0.533","✗"],
    ["ResNet-50","image","0.572","✗"],
    ["SO-SDC-Prioritizer","GA (TOSEM’23)","0.765","✗"],
    ["ITEP4SDC","MLP (ICST’25)","0.781","✗"],
    ["Greedy-diversity","heuristic","0.795","✗"],
    ["RoadFury","Transformer+SWA","0.804","✗"],
    ["SE2RoadNet","SE(2)-equiv.","0.805","✓ Δ=0"]],
    [2.3,2.1,1.0,1.0], fontsize=13, hi=9, row_h=0.42, aligns=['l','l','r','c'])
# native 2-axis scatter (APFD x  vs  invariance y)
text(s, 7.1, 1.25, 5.7, 0.4, [[B("Hai trục độc lập", TEAL)]], size=15, align='c')
ax_x, ax_y, ax_w, ax_h = 8.0, 5.1, 4.3, 3.3
axis(s, ax_x, ax_y-ax_h, ax_w, ax_h, NAVYMID)
text(s, ax_x+ax_w-0.6, ax_y+0.05, 0.9, 0.3, [[{'t':"APFD",'sz':11,'c':TEAL}]])
text(s, ax_x-0.5, ax_y-ax_h-0.05, 1.6, 0.3, [[{'t':"bất biến?",'sz':11,'c':TEAL}]])
base_y = ax_y - 0.25
for fx in (0.6, 0.95, 1.35, 2.6, 2.9, 3.2, 3.55):   # prior methods on baseline
    d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(ax_x+fx-0.05), Inches(base_y-0.05), Inches(0.1), Inches(0.1))
    d.fill.solid(); d.fill.fore_color.rgb=NAVYMID; d.line.fill.background(); d.shadow.inherit=False
text(s, ax_x+0.2, base_y-0.55, 3.6, 0.3, [[{'t':"prior: Δ > 0 (phụ thuộc khung)",'sz':10,'c':NAVYMID}]])
se = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(ax_x+3.5-0.08), Inches(ax_y-ax_h+0.25-0.08), Inches(0.16), Inches(0.16))
se.fill.solid(); se.fill.fore_color.rgb=ORANGE; se.line.fill.background(); se.shadow.inherit=False
text(s, ax_x+2.6, ax_y-ax_h-0.02, 1.7, 0.3, [[{'t':"SE2RoadNet",'b':True,'c':ORANGE,'sz':10}]])
arrow(s, ax_x+3.5, base_y-0.15, ax_x+3.5, ax_y-ax_h+0.45, ORANGE, 1.6)
callout(s, 0.7, 6.25, 11.9, 0.85,
        [[B("Khoảng trống: ",TEAL),I("cùng một con đường — chỉ khác hướng đặt hoặc số điểm lấy mẫu — không được làm điểm số thay đổi. Chưa method nào đảm bảo điều đó.")]], size=14)

# ============ 12. divider ============
divider("RoadFury: Phương pháp nền tảng")

# ============ 13. RoadFury pipeline ============
s = blank(prs); bar(s, "RoadFury: Pipeline nền tảng", "10/37")
imgfit(s, A+"fig_arch-1.png", 0.4, 1.05, 12.5, 6.2, valign='t')

# ============ 14. RoadFury detail ============
s = blank(prs); bar(s, "RoadFury (chi tiết): mạnh, nhưng có điểm mù", "11/37")
text(s, 0.5, 1.1, 6.2, 0.35, [[B("Trích xuất đặc trưng (10 kênh)")]], size=15)
feat = [["f₀","Chiều dài đoạn ‖pᵢ₊₁−pᵢ‖"],["f₁","Biến thiên góc tuyệt đối |Δθᵢ|"],
        ["f₂","Độ cong Menger κᵢ"],["f₃","Biến thiên độ cong Δκᵢ (jerk)"],
        ["f₄","Chiều dài cung tích luỹ s/L"],["f₅","sin θᵢ  (phụ thuộc khung)"],
        ["f₆","cos θᵢ  (phụ thuộc khung)"],["f₇","Hướng tuyệt đối θᵢ  (phụ thuộc khung)"],
        ["f₈","Độ lệch chuẩn cục bộ σ(κ)"],["f₉","Gia tốc cong Δ²κᵢ"]]
gt = table(s, 0.5, 1.5, 6.2, [["Kênh","Ý nghĩa"]]+feat, [0.9,5.3], fontsize=12, header=True, row_h=0.345, aligns=['c','l'])
for ri in (6,7,8):  # highlight frame-dependent rows (f5,f6,f7 -> data rows 6,7,8)
    for ci in (0,1):
        gt.cell(ri,ci).fill.solid(); gt.cell(ri,ci).fill.fore_color.rgb = RGBColor(0xF6,0xE7,0xE6)
        gt.cell(ri,ci).text_frame.paragraphs[0].runs[0].font.color.rgb = ALERT
text(s, 7.0, 1.1, 5.8, 0.35, [[B("RoadTransformer (829K tham số)")]], size=15)
bullets(s, 7.0, 1.5, 5.8, 2.3, [
    "Linear 10→128 + LayerNorm + GELU.",
    [I("Token [CLS] + learned PE "),B("(vị trí tuyệt đối)",ALERT),I(".")],
    "4 lớp Pre-LN Transformer, 8 heads, FFN 512.",
    "Pool [CLS] → 128→64→1 → Sigmoid."], size=13.5, gap=6)
card(s, 7.0, 3.95, 5.8, 1.15, "Kết quả ICST",
     [[I("APFD = "),B("0.804 ± 0.012",GREEN),I(", AUC 0.917 (30 trials, 956 test). Tốt nhất competition.")]],
     RGBColor(0xEAF1EA & 0xFFFFFF if False else 0xEA,0xF1,0xEA), GREEN, tcolor=GREEN, size=13.5)
callout(s, 0.7, 6.35, 11.9, 0.8,
        [[I("3/10 kênh và PE tuyệt đối phụ thuộc khung quy chiếu ⇒ APFD cao nhưng "),B("không có bảo chứng",ALERT),I(" — đây là 2 điểm mù.")]], size=14)

# ============ 15. bridge (NATIVE) ============
s = blank(prs); bar(s, "Từ RoadFury đến SE2RoadNet: Vá 2 điểm mù", "12/37")
GAPF = RGBColor(0xFB, 0xEC, 0xEB)
def gapcard(x, num, title, body):
    box(s, x, 1.2, 4.5, 1.45, GAPF, ALERT)
    badge(s, x+0.35, 1.55, 0.4, num, ALERT)
    text(s, x+0.2, 1.32, 4.1, 1.3, [[B("✗ "+title, ALERT)], [{'t':body,'i':True,'sz':11}]], size=14, align='c')
def fixcard(x, title, body):
    box(s, x, 3.35, 4.5, 1.2, GTINT, GREEN)
    text(s, x+0.2, 3.5, 4.1, 1.0, [[B("✓ "+title, GREEN)], [{'t':body,'i':True,'sz':11}]], size=14, align='c')
gapcard(1.6, 1, "Nhạy phép xoay", "cùng đường, xoay 30°/90° ⇒ điểm trôi (ITEP4SDC Δ=0.057)")
gapcard(7.2, 2, "Nhạy tần số lấy mẫu", "cùng đường, N=64 vs 197 điểm ⇒ điểm đổi theo cách rời rạc hoá")
fixcard(1.6, "7 kênh bất biến SE(2)", "bỏ sin θ, cos θ, θ; chỉ giữ Δs, κ, dκ/ds…")
fixcard(7.2, "Attention bias theo Δs", "thay PE tuyệt đối bằng hiệu số cung — bất biến")
for x in (3.85, 9.45):
    arrow(s, x, 2.7, x, 3.3, ORANGE, 2.4)
    chip(s, x-0.35, 2.82, 0.7, 0.32, [[{'t':"vá",'b':True,'c':WHITE,'sz':11}]], line=ORANGE, fill=ORANGE, color=WHITE)
box(s, 1.15, 4.95, 11.0, 1.15, NTINT, TEAL, width=1.6)
text(s, 1.3, 5.05, 10.7, 1.0, [
    [B("SE2RoadNet: Bất biến SE(2) bằng kiến trúc", TEAL)],
    [B("Augmentation (xấp xỉ)", ALERT), I("   ──(paradigm shift)──▶   "), B("ràng buộc kiến trúc (bảo chứng)", GREEN)]],
    size=15, align='c', anchor='m', sp_after=4)
text(s, 0, 6.35, W, 0.4, [[I("Mỗi điểm mù RoadFury → một thành phần SE2RoadNet: "),B("vá có chủ đích",TEAL),I(", không refactor mù.")]], size=13.5, align='c')

# ============ 16. divider ============
divider("SE2RoadNet: Đề xuất của nhóm")

# ============ 17. SE2 pipeline ============
s = blank(prs); bar(s, "SE2RoadNet: Tổng quan kiến trúc", "13/37")
imgfit(s, A+"fig_arch_se2-1.png", 0.4, 1.05, 12.5, 6.2, valign='t')

# ============ 18. SE(2) core ============
s = blank(prs); bar(s, "Ý tưởng cốt lõi: Bất biến SE(2)", "14/37")
text(s, 0.6, 1.1, 12.1, 0.7, [
    [B("Định lý mục tiêu:  ",TEAL),I("Mô hình f"),{'t':"θ",'sz':10},I(" phải thoả mãn, "),B("bit-identical"),I(", với mọi (R, t) ∈ SE(2):")]], size=15)
imgfit(s, A+"fig-15.png", 4.3, 1.6, 4.7, 0.85, valign='t')
text(s, 0.6, 2.55, 12.1, 0.4, [[{'t':"Xoay, tịnh tiến, hoặc đổi gốc toạ độ con đường ⇒ điểm số y nguyên.",'i':True}]], size=15, align='c')
def se2panel(px, py, label, rot=0.0, shift=0.0, reframe=False):
    box(s, px, py, 3.7, 1.6, SOFT, TINT)
    text(s, px, py+0.05, 3.7, 0.3, [[B(label, TEAL)]], size=12.5, align='c')
    axis(s, px+0.5, py+0.5, 2.5, 0.55, NAVYMID)
    sine(s, px+0.6+shift, py+0.42, 2.3, 0.58, ORANGE, 2.4, cycles=1.3, amp=0.7, rot=rot, slope=0.12)
    text(s, px, py+1.18, 3.7, 0.3, [[{'t':"Score = 0.8047",'b':True,'c':GREEN,'sz':11}]], align='c')
se2panel(2.7, 3.05, "Gốc R")
se2panel(7.0, 3.05, "Xoay R", rot=40)
se2panel(2.7, 4.85, "Tịnh tiến +t", shift=0.25)
se2panel(7.0, 4.85, "Đổi gốc toạ độ", reframe=True)
for ey in (3.82, 5.62):
    text(s, 6.55, ey, 0.6, 0.4, [[{'t':"=",'b':True,'c':TEAL,'sz':22}]], align='c')
text(s, 5.3, 4.55, 0.6, 0.4, [[{'t':"=",'b':True,'c':TEAL,'sz':22}]], align='c')

# ============ 19. 7 channels ============
s = blank(prs); bar(s, "Bước 1: 7 kênh đặc trưng bất biến SE(2)", "15/37")
text(s, 0.5, 1.1, 7.0, 0.35, [[B("Chỉ giữ đại lượng nội tại của đường")]], size=15)
seven = ["Δsᵢ – chiều dài đoạn i",
         "|Δθᵢ| – biến thiên góc tuyệt đối",
         "κᵢ = Δθᵢ/Δsᵢ – độ cong có dấu",
         "dκ/ds – đạo hàm độ cong (jerk hình học)",
         "d²κ/ds² – đạo hàm bậc 2",
         "s_norm = s/L ∈ [0,1] – chiều dài cung tương đối",
         "σ_local(κ) – độ lệch chuẩn cục bộ (cửa sổ 11)"]
tb = s.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(7.0), Inches(3.9)); tf=tb.text_frame; tf.word_wrap=True
for i,it in enumerate(seven):
    p = tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(9)
    rn=p.add_run(); rn.text=f"{i+1}.  "; rn.font.bold=True; rn.font.color.rgb=ORANGE; rn.font.size=Pt(14); rn.font.name="Calibri"
    r=p.add_run(); r.text=it; r.font.size=Pt(14); r.font.name="Calibri"; r.font.color.rgb=TEAL
# native road sketch with intrinsic quantities
box(s, 7.9, 1.5, 5.0, 2.8, WHITE, TINT)
sine(s, 8.2, 1.9, 4.4, 1.8, ORANGE, 3.0, cycles=1.5, amp=0.7, slope=0.0)
import math as _m
for tt in (0.30, 0.45, 0.60):
    rx = 8.2 + tt*4.4; ry = 1.9 + 0.9 - 0.7*0.9*_m.sin(2*_m.pi*1.5*tt)
    dd = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(rx-0.05), Inches(ry-0.05), Inches(0.1), Inches(0.1))
    dd.fill.solid(); dd.fill.fore_color.rgb=TEAL; dd.line.fill.background(); dd.shadow.inherit=False
text(s, 8.2, 3.6, 4.4, 0.4, [[{'t':"Δs,  κ,  dκ/ds,  …",'b':True,'c':TEAL,'sz':12}]], align='c')
text(s, 8.0, 4.35, 4.9, 0.9, [[{'t':"Đại lượng nội tại tính trên từng cặp điểm liền kề — không đổi khi xoay/tịnh tiến.",'i':True,'c':GREY}]], size=13.5, align='c')
callout(s, 1.5, 6.3, 10.3, 0.85, [[B("⇒ 0 kênh phụ thuộc khung quy chiếu",ORANGE),I("  (đã bỏ sin θ, cos θ, θ).")]], size=15)

# ============ 20. architecture stack (NATIVE) ============
s = blank(prs); bar(s, "Bước 2: Kiến trúc SE2RoadNet (InvariantBlock ×6)", "16/37")
stk = [("Linear 7→192 + LN + GELU", TINT, TEAL),
       ("Concat token [CLS]", TINT, TEAL),
       ("InvariantBlock ×6", OTINT, ORANGE),
       ("Pool theo [CLS]", TINT, TEAL),
       ("Head: LN → 64 → 1", GTINT, GREEN)]
sx, sw, sh0 = 1.3, 4.0, 0.72
for i, (label, fill, line) in enumerate(stk):
    yy = 1.45 + i*1.02
    box(s, sx, yy, sw, sh0, fill, line, width=1.4)
    text(s, sx, yy, sw, sh0, [[B(label, line if line!=TEAL else TEAL)]], size=15, align='c', anchor='m')
    if i < 4:
        arrow(s, sx+sw/2, yy+sh0, sx+sw/2, yy+1.02, ORANGE, 2.2)
bullets(s, 6.6, 1.7, 6.3, 3.0, [
    [I("Mỗi điểm trên đường ⇒ một "),B("token 192 chiều"),I(".")],
    [B("Mỗi block"),I(" = MHA (8 heads) + FFN (512) + LayerNorm, với "),B("relative-arclength bias",TEAL),I(" (Bước 3).")],
    [I("Pool [CLS] → 1 logit = xác suất FAIL.")]], size=15, gap=11)
card(s, 6.6, 5.0, 6.3, 1.3, "Quy mô",
     [[I("Tổng "),B("2.11 M",TEAL),I(" tham số (~2.5× RoadFury) — vẫn nhẹ, train 24 phút.")]],
     RGBColor(0xEA,0xF1,0xEA), GREEN, tcolor=GREEN, size=14)

# ============ 21. attention ============
s = blank(prs); bar(s, "Bước 3: Attention với bias hiệu số Δs", "17/37")
text(s, 0.5, 1.1, 6.4, 0.35, [[B("Vấn đề với PE tuyệt đối")]], size=15)
text(s, 0.5, 1.5, 6.4, 1.1, [[I("PE chuẩn (sin/cos của index i) khoá mô hình vào "),B("vị trí tuyệt đối"),I(" của token. Xoay/cắt đường khác ⇒ index ứng arclength khác ⇒ vỡ bất biến.")]], size=14)
text(s, 0.5, 2.75, 6.4, 0.35, [[B("Giải pháp: relative-arclength RFF bias")]], size=15)
text(s, 0.5, 3.15, 6.4, 0.9, [[I("Thêm bias vào ma trận attention, chỉ phụ thuộc "),B("hiệu số"),I(" Δsᵢⱼ = sᵢ − sⱼ (bất biến tái tham số hoá):")]], size=14)
imgfit(s, A+"fig-16.png", 0.7, 4.1, 5.7, 0.8, valign='t')
text(s, 0.5, 5.0, 6.4, 0.5, [[I("với ω ∈ ℝ"),{'t':"32",'sz':9},I(" là tần số Fourier ngẫu nhiên (RFF) cố định.")]], size=14)
# native tokens + Δs arrow + heatmap
def _blend(a):  # a in [0,1] -> orange over white
    return RGBColor(int(255*(1-a)+0xEB*a), int(255*(1-a)+0x81*a), int(255*(1-a)+0x1B*a))
tok_x, tok_y, tw, gap = 7.4, 1.3, 0.62, 0.16
for i in range(6):
    x = tok_x + i*(tw+gap)
    hl = i in (2, 4)
    box(s, x, tok_y, tw, 0.55, _blend(0.3) if hl else SOFT, ORANGE if hl else TEAL)
    text(s, x, tok_y, tw, 0.55, [[{'t':f"s{i+1}",'b':True,'c':TEAL,'sz':13}]], align='c', anchor='m')
sx3 = tok_x + 2*(tw+gap) + tw/2; sx5 = tok_x + 4*(tw+gap) + tw/2
arrow(s, sx3, tok_y+0.7, sx5, tok_y+0.7, ORANGE, 1.8)
text(s, (sx3+sx5)/2-1.0, tok_y+0.78, 2.0, 0.3, [[{'t':"Δs = s₅ − s₃",'c':ORANGE,'sz':12}]], align='c')
import math as _mm
hm_x, hm_y, cell = 8.55, 2.55, 0.36
for r_ in range(6):
    for c_ in range(6):
        a = _mm.exp(-0.55*abs(r_-c_))
        cb = box(s, hm_x+c_*cell, hm_y+r_*cell, cell, cell, _blend(a*0.85), WHITE, rounded=False, width=0.5)
text(s, 6.9, 5.0, 6.0, 0.6, [[{'t':"Heatmap: attention giảm theo |Δsᵢⱼ| — bias chỉ phụ thuộc hiệu số.",'i':True,'c':GREY}]], size=12.5, align='c')
callout(s, 0.7, 6.35, 11.9, 0.8, [[B("Chi phí: ",TEAL),I("O(B·L²·32) mỗi block — điểm tốn nhất (24 phút train). Cải tiến tương lai: RoPE-1D trên s_norm.")]], size=13.5)

# ============ 22. training ============
s = blank(prs); bar(s, "Bước 4: Huấn luyện", "18/37")
imgfit(s, F+"gen_focal_loss.png", 0.5, 1.2, 6.5, 4.6, valign='t')
text(s, 7.3, 1.2, 5.6, 0.35, [[B("Vì sao Focal Loss?")]], size=15)
imgfit(s, A+"fig-17.png", 7.5, 1.65, 5.0, 0.55, valign='t')
text(s, 7.3, 2.35, 5.6, 1.1, [[I("FAIL chỉ ~30–40%; BCE thuần ⇒ model “đoán PASS hết”. Hệ số (1−p̂ₜ)^γ "),B("down-weight",TEAL),I(" ca dễ, "),B("up-weight",TEAL),I(" ca khó.")]], size=14)
text(s, 7.3, 3.55, 5.6, 0.35, [[B("Cấu hình")]], size=15)
table(s, 7.3, 3.95, 5.6, [
    ["Optimizer","AdamW (wd 1e-3)"],["LR","5e-4 cosine + warmup 5"],
    ["Batch / Epochs","384 / 80 (SWA từ 56)"],["Precision","bf16"],
    ["Sampler","WeightedRandom"],["Wall-clock / Params","24.2 phút  ·  2.11 M"]],
    [2.2,3.4], fontsize=13, header=False, row_h=0.36, aligns=['l','l'])

# ============ 23. divider ============
divider("Đánh giá thực nghiệm")

# ============ 24. APFD ============
s = blank(prs); bar(s, "Chỉ số đánh giá: APFD", "19/37")
text(s, 0.5, 1.15, 6.4, 0.7, [[B("Average Percentage of Faults Detected")],
     [I("Với hoán vị π trên n test (m là FAIL), TFᵢ là vị trí fault thứ i:")]], size=15, sp_after=4)
imgfit(s, A+"fig-18.png", 0.7, 2.35, 5.8, 1.2, valign='t')
bullets(s, 0.5, 3.9, 6.4, 1.6, [
    [B("Khoảng: "),I("[0, 1], càng cao càng tốt.")],
    [B("Ngẫu nhiên",GREY),I(" ≈ 0.5;  "),B("Lý tưởng",GREEN),I(" ≈ 1.0.")]], size=15, gap=8)
imgfit(s, F+"gen_apfd_curve.png", 7.0, 1.4, 5.9, 4.4, valign='t')
text(s, 7.0, 5.85, 5.9, 0.4, [[{'t':"Ranking tốt ⇒ đường cong tăng dốc đứng từ đầu.",'i':True,'c':GREY}]], size=13.5, align='c')

# ============ 25. eval protocol (NATIVE) ============
s = blank(prs); bar(s, "Giao thức đánh giá", "20/37")
panels = [
    (0.55, "1. Single-pass", TEAL, SOFT, TEAL,
     ["Xếp hạng 1 lần trên toàn", "Competition split (956 test)"], "APFD = 0.8047", None),
    (4.72, "2. Multi-trial (30 lần)", ORANGE, OTINT, ORANGE,
     ["Mỗi trial: lấy ngẫu nhiên", "287/956 test → APFD"], "0.8048 ± 0.0118", "loại “may rủi” do thứ tự cố định"),
    (8.89, "3. Rotation probe", GREEN, GTINT, GREEN,
     ["Xoay toàn bộ split bằng", "{0, 30, 60, 90, 180, −45}°", "rồi đánh giá lại"], "Δ = max − min APFD", None),
]
pw = 3.85
for x, ttl, c, fill, line, body, res, note in panels:
    box(s, x, 1.3, pw, 4.3, fill, line)
    text(s, x, 1.55, pw, 0.4, [[B(ttl, c)]], size=16, align='c')
    rule(s, x+0.5, 2.1, pw-1.0, line)
    text(s, x+0.2, 2.35, pw-0.4, 1.2, [[I(l)] for l in body], size=15, align='c')
    chip(s, x+0.55, 3.75, pw-1.1, 0.55, [[B(res, c)]], line=line, color=TEAL, size=15)
    if note:
        text(s, x+0.2, 4.5, pw-0.4, 0.4, [[{'t':note,'i':True,'c':GREY,'sz':11}]], align='c')
arrow(s, 4.42, 3.45, 4.70, 3.45, NAVYMID, 2.4)
arrow(s, 8.59, 3.45, 8.87, 3.45, NAVYMID, 2.4)
callout(s, 1.0, 6.3, 11.3, 0.8, [[B("Ba lớp: ",TEAL),I("single-pass (điểm chính) → multi-trial (độ ổn định) → rotation probe (kiểm chứng bất biến).")]], size=15)

# ============ 26. headline ============
s = blank(prs); bar(s, "Kết quả headline: Bất biến phép xoay Δ = 0", "21/37")
imgfit(s, F+"gen_rotation_drift.png", 0.5, 1.2, 7.0, 4.5, valign='t')
table(s, 7.8, 1.2, 5.0, [
    ["Xoay","ITEP4SDC","SE2RoadNet"],
    ["0°","0.7810","0.8047"],["+30°","0.7240","0.8047"],["+60°","0.7518","0.8047"],
    ["+90°","0.7396","0.8047"],["+180°","0.7334","0.8047"],["−45°","0.7627","0.8047"],
    ["Δ","0.0570","0.0000"]],
    [1.4,1.8,1.8], fontsize=13.5, hi=7, aligns=['l','c','c'])
callout(s, 7.8, 4.55, 5.0, 1.35,
        [[B("Không phải “trong sai số”:")],[I("6 góc xoay, APFD bằng nhau đến từng bit float (= 0) ⇒ pipeline 7 kênh trả vector pixel-identical sau xoay.")]],
        line=GREEN, size=13.5, align='l')
text(s, 0, 6.5, W, 0.4, [[B("⇒ Lý thuyết được xác minh bằng thực nghiệm",ORANGE),I(", không cần nói “xấp xỉ”.")]], size=15, align='c')

# ============ 27. leaderboard ============
s = blank(prs); bar(s, "Bảng xếp hạng: so với 8 baselines", "22/37")
imgfit(s, F+"gen_leaderboard.png", 0.5, 1.2, 7.2, 4.5, valign='t')
table(s, 8.0, 1.4, 4.7, [
    ["Method","AUC","APFD"],
    ["RoadFury","0.917","0.804"],
    ["SE2RoadNet","0.934","0.805"]],
    [2.5,1.1,1.1], fontsize=14, hi=2, aligns=['l','c','c'])
text(s, 8.0, 2.7, 4.7, 0.35, [[B("Đọc bảng",GREEN)]], size=15)
text(s, 8.0, 3.1, 4.7, 2.0, [[I("SE2RoadNet "),B("ngang APFD"),I(" baseline tốt nhất, nhưng "),B("tăng AUC"),I(" (+0.017) và "),B("thêm bảo chứng"),I(" Δ=0 mà không method nào có ⇒ cải thiện theo nghĩa "),{'t':"Pareto",'i':True},I(".")]], size=14)
callout(s, 0.7, 6.35, 11.9, 0.8, ["Không đánh đổi: giữ APFD đỉnh, tăng AUC, và là phương pháp duy nhất bất biến SE(2) tuyệt đối."], size=15)

# ============ 28. divider ============
divider("Hạn chế & Hướng phát triển")

# ============ 29. limitations ============
s = blank(prs); bar(s, "Hạn chế: nhìn thẳng vào điểm yếu", "23/37")
text(s, 0.6, 1.15, 6.0, 0.35, [[B("Về chỉ số",TEAL)]], size=15)
bullets(s, 0.6, 1.6, 6.0, 4.0, [
    [B("AUC và APFD phân kỳ: "),I("AUC cao hơn không kéo theo APFD cao hơn — phải lập luận cẩn thận metric nào quan trọng khi nào.")],
    [B("Listwise loss chưa tăng APFD trung bình: "),I("chỉ giảm σ ⇒ đóng góp “ổn định”, không phải headline.")],
    [B("Trần APFD ở bench FAIL cao: "),I("ví dụ 95% FAIL ⇒ APFD ≈ 0.52 là ceiling, không phải thất bại.")]], size=14, gap=10, color=ALERT)
text(s, 6.9, 1.15, 6.0, 0.35, [[B("Về phương pháp",TEAL)]], size=15)
bullets(s, 6.9, 1.6, 6.0, 4.0, [
    [B("Conformal an toàn còn dở: "),I("v1 valid-nhưng-vô nghĩa, v2 informative-nhưng-invalid; cần v3 mới ship.")],
    [B("Dịch phân phối chưa đóng: "),I("IRM / TENT chưa thu hẹp gap SensoDat→Competition (negative đã biết).")],
    [B("Bất biến lấy mẫu là xấp xỉ: "),I("resolution probe Δ ≈ 0.0012 (rất nhỏ) — chưa exact như phép xoay.")]], size=14, gap=10, color=ALERT)
callout(s, 0.7, 6.35, 11.9, 0.8,
        ["Bất biến xoay là exact; các đảm bảo khác (resolution, conformal) vẫn đang hoàn thiện — báo cáo trung thực thay vì giấu."], size=14)

# ============ 30. future 4-cards (NATIVE) ============
s = blank(prs); bar(s, "Hướng phát triển", "24/37")
fut = [(0.6, 1.25, "1", "Tổng quát hoá đa benchmark", "Một công thức, 8+ benchmark công khai (OOB, SDC-Scissor, sdc-travel…) — APFD không tinh chỉnh theo từng bench."),
       (6.9, 1.25, "2", "An toàn có bảo chứng (Conformal v3)", "Chặn dưới prefix-APFD vừa valid vừa non-vacuous; audit tỉ lệ vi phạm độ cong."),
       (0.6, 3.75, "3", "Bất biến lấy mẫu exact", "RoPE-1D trên s_norm thay RFF bias ⇒ rẻ hơn và tiến tới resolution-Δ = 0."),
       (6.9, 3.75, "4", "SSL vật lý (physics-informed)", "Pretext task gắn với độ cong/jerk để học biểu diễn nền — thay SSL hình học ngây thơ (chuyển kém).")]
for x, y, num, title, body in fut:
    box(s, x, y, 5.85, 2.25, OTINT, ORANGE)
    text(s, x+5.0, y+0.05, 0.7, 0.7, [[{'t':num,'b':True,'c':RGBColor(0xF3,0xCF,0xA8),'sz':40}]], align='r')
    text(s, x+0.25, y+0.25, 5.0, 0.6, [[B(title, TEAL)]], size=15)
    text(s, x+0.25, y+1.0, 5.35, 1.1, [[I(body)]], size=13.5)
text(s, 0, 6.5, W, 0.4, [[I("⇒ Mục tiêu: từ một baseline bất biến cho SDC → "),B("công thức chung",TEAL),I(" bất biến & audit-readable cho mọi benchmark.")]], size=14, align='c')

# ============ 31. references ============
s = blank(prs); bar(s, "Tài liệu tham khảo", "25/37")
refs = [
    [B("1. "),I("C. Birchler, C. Rohrbach, T. Kehrer, S. Panichella. "),{'t':"SensoDat: Simulation-based Sensor Dataset of Self-driving Cars",'i':True},I(". MSR 2024.")],
    [B("2. "),I("C. Birchler et al. "),{'t':"Single and Multi-objective Test Cases Prioritization for Self-driving Cars in Virtual Environments",'i':True},I(". ACM TOSEM 2023.")],
    [B("3. "),{'t':"ICST/SBFT Tool Competition — Self-Driving Car Testing (Cyber-Physical Systems)",'i':True},I(". 2023–2025.")],
    [B("4. "),I("G. Rothermel, R. Untch, C. Chu, M. Harrold. "),{'t':"Prioritizing Test Cases for Regression Testing",'i':True},I(". IEEE TSE 2001.")],
    [B("5. "),I("A. Vaswani et al. "),{'t':"Attention Is All You Need",'i':True},I(". NeurIPS 2017.")],
    [B("6. "),I("T.-Y. Lin et al. "),{'t':"Focal Loss for Dense Object Detection",'i':True},I(". ICCV 2017.")],
    [B("7. "),I("P. Izmailov et al. "),{'t':"Averaging Weights Leads to Wider Optima and Better Generalization (SWA)",'i':True},I(". UAI 2018.")],
    [B("8. "),I("T. Cohen, M. Welling. "),{'t':"Group Equivariant Convolutional Networks",'i':True},I(". ICML 2016.")],
    [B("9. "),I("A. Rahimi, B. Recht. "),{'t':"Random Features for Large-Scale Kernel Machines (RFF)",'i':True},I(". NeurIPS 2007.")],
    [B("10. "),I("A. Angelopoulos, S. Bates. "),{'t':"A Gentle Introduction to Conformal Prediction",'i':True},I(". 2021.")],
]
text(s, 0.7, 1.2, 12.0, 5.6, refs, size=13.5, sp_after=7)

# ============ 32. thanks ============
s = blank(prs)
text(s, 0, 1.7, W, 1.0, [[B("Cảm ơn thầy đã lắng nghe!",TEAL)]], size=34, align='c')
text(s, 0, 2.9, W, 0.6, [[{'t':"Q&A — Questions & Discussion",'c':ORANGE}]], size=20, align='c')
imgfit(s, F+"hcmus.png", 1.3, 4.6, 1.5, 1.5)
imgfit(s, F+"icst_logo.jpg", W-2.7, 4.7, 1.4, 1.3)
text(s, 0, 4.7, W, 1.4, [
    [B("Test Prioritization for Self-Driving Cars",TEAL)],
    [I("SE2RoadNet — bất biến SE(2) bằng kiến trúc")],
    [I("Trần Chí Nguyên · Huỳnh Trung Kiệt · Đào Sỹ Duy Minh")],
    [I("University of Science, VNU-HCM · ICST 2026")]], size=15, align='c', sp_after=6)

# ============ 33. backup APFD example ============
s = blank(prs); bar(s, "Backup: APFD — ví dụ tính tay", "26/37")
text(s, 0.6, 1.2, 12.0, 0.4, [[I("Giả sử n = 5 test, m = 2 FAIL. Hai cách xếp hạng:")]], size=15)
text(s, 1.4, 2.0, 5.0, 0.4, [[B("Xếp tốt",GREEN),I("  (FAIL ở vị trí 1, 2)")]], size=15, align='c')
imgfit(s, A+"fig-19.png", 1.4, 2.5, 5.0, 1.0, valign='t')
text(s, 7.0, 2.0, 5.0, 0.4, [[B("Xếp kém",ALERT),I("  (FAIL ở vị trí 4, 5)")]], size=15, align='c')
imgfit(s, A+"fig-20.png", 7.0, 2.5, 5.0, 1.0, valign='t')
callout(s, 1.2, 5.0, 10.9, 1.0,
        ["Cùng một tập lỗi — chỉ khác thứ tự — APFD chênh nhau 4 lần. Đó là toàn bộ giá trị của test prioritization."], size=15)

# ============ 34. backup multi-trial ============
s = blank(prs); bar(s, "Backup: Multi-trial protocol (vì sao 30 lần?)", "27/37")
bullets(s, 0.7, 1.4, 11.9, 3.6, [
    "APFD single-pass nhạy với thứ tự cố định của tập test ⇒ một con số dễ “may rủi”.",
    "Mỗi trial: lấy ngẫu nhiên max(50, 0.3·|test|) = 287/956 test, tính APFD độc lập.",
    "Báo cáo APFD ± σ qua 30 trials, seed = 42.",
    "σ thường là con số xuất bản quan trọng hơn mean (đo độ ổn định của ranking)."], size=16, gap=12)
callout(s, 3.0, 5.3, 7.3, 0.9, [[I("SE2RoadNet: "),B("0.8048 ± 0.0118",TEAL),I("  (mean ± σ, 30 trials)")]], size=15)

# ============ 35. backup AUC vs APFD ============
s = blank(prs); bar(s, "Backup: AUC vs APFD — vì sao tách bạch?", "28/37")
card(s, 0.7, 1.3, 5.9, 1.3, "AUC", [[I("đo khả năng phân loại PASS/FAIL trên từng cặp — không quan tâm thứ hạng tuyệt đối.")]], SOFT, TEAL, size=14)
card(s, 6.8, 1.3, 5.9, 1.3, "APFD", [[I("đo tốc độ lộ lỗi theo thứ tự chạy — chính là cái competition tối ưu.")]], SOFT, ORANGE, tcolor=ORANGE, size=14)
bullets(s, 0.7, 3.0, 11.9, 2.6, [
    "Trong gần như mọi thí nghiệm của nhóm, AUC và APFD phân kỳ: tăng AUC không đảm bảo tăng APFD.",
    "SE2RoadNet tăng AUC (0.917 → 0.934) trong khi APFD giữ đỉnh (≈ 0.805).",
    "⇒ Khi báo cáo phải nói rõ metric nào đang được tối ưu cho tình huống nào."], size=15, gap=11)
callout(s, 0.7, 6.0, 11.9, 0.85, ["Không gộp AUC và APFD thành “một điểm tốt hơn” — chúng đo hai thứ khác nhau."], size=15)

# ============ 36. backup resolution ============
s = blank(prs); bar(s, "Backup: Resolution probe (bất biến lấy mẫu)", "29/37")
bullets(s, 0.7, 1.4, 11.9, 3.5, [
    "Lấy cùng một con đường, resample ở N ∈ {64, …, 197} điểm rồi đánh giá lại.",
    "Đo Δ = max − min APFD qua các mức N.",
    "Kết quả: Δ ≈ 0.0012 — rất nhỏ nhưng chưa exact như phép xoay (Δ = 0).",
    "Lý do: bias RFF theo Δs là gần-bất-biến tái tham số hoá, không tuyệt đối ⇒ hướng RoPE-1D."], size=16, gap=12)
callout(s, 0.9, 5.5, 11.5, 0.9, ["Phép xoay: Δ = 0 exact (bằng kiến trúc). Lấy mẫu: Δ ≈ 0.0012 (gần như, đang tiến tới exact)."], size=15)

# ============ 37. backup cost ============
s = blank(prs); bar(s, "Backup: Chi phí ở quy mô lớn", "30/37")
bullets(s, 0.7, 1.4, 11.9, 3.5, [
    "Inference: 1 forward pass/đường, batch (n, 7, 197) ⇒ chấm cả tập trong vài giây GPU.",
    "Chi phí thật nằm ở mô phỏng vật lý của các test được chọn — không phải ở scorer.",
    "Ranking tốt cắt 50–80% số kịch bản cần chạy mô phỏng ⇒ tiết kiệm nhiều giờ CPU.",
    "Train một lần: 24.2 phút (Kaggle), 2.11 M tham số ⇒ tái lập rẻ."], size=16, gap=12)
callout(s, 0.7, 5.5, 11.9, 0.9, [[{'t':"“The best simulation is the one you never run.”",'i':True,'b':True},I(" — scorer rẻ, ngân sách dồn cho các test có khả năng lộ lỗi cao nhất.")]], size=15)

prs.save("SE2RoadNet_native_editable.pptx")
print("saved SE2RoadNet_native_editable.pptx |", len(prs.slides._sldIdLst), "slides")
