# -*- coding: utf-8 -*-
"""
Clean the LibreOffice-converted editable pptx:
  1. Merge split title fragments (white FAFAFA boxes at top) into one clean title.
  2. Merge section-divider title fragments.
  3. Add a space between touching NORMAL-HEIGHT prose boxes (fix glued words),
     while NEVER touching short sub/superscript boxes (math stays tight).
Text content is otherwise preserved -> formulas (fraction layouts) untouched.
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

SRC = "SE2RoadNet_editable_raw_backup.pptx"
DST = "SE2RoadNet_Test_Prioritization_SDC_editable.pptx"

# clean titles (1-based slide index). Dividers included.
TITLES = {
    2:  "Nội dung chính",
    3:  "Bối cảnh & Phát biểu bài toán",
    4:  "Bối cảnh: Kiểm thử Self-Driving Car trong mô phỏng",
    5:  "Phát biểu bài toán & Tiêu chí đánh giá",
    6:  "Vì sao cần ưu tiên kiểm thử?",
    8:  "Ví dụ kịch bản: hình dạng quyết định nhãn",
    9:  "Công trình liên quan & Khoảng trống nghiên cứu",
    10: "Related Work: 3 hướng tiếp cận",
    11: "Related Work: Tổng hợp & Khoảng trống",
    12: "RoadFury: Phương pháp nền tảng",
    13: "RoadFury: Pipeline nền tảng",
    14: "RoadFury (chi tiết): mạnh, nhưng có điểm mù",
    15: "Từ RoadFury đến SE2RoadNet: Vá 2 điểm mù",
    16: "SE2RoadNet: Đề xuất của nhóm",
    17: "SE2RoadNet: Tổng quan kiến trúc",
    18: "Ý tưởng cốt lõi: Bất biến SE(2)",
    19: "Bước 1: 7 kênh đặc trưng bất biến SE(2)",
    20: "Bước 2: Kiến trúc SE2RoadNet (InvariantBlock ×6)",
    21: "Bước 3: Attention với bias hiệu số ∆s",
    22: "Bước 4: Huấn luyện",
    23: "Đánh giá thực nghiệm",
    24: "Chỉ số đánh giá: APFD",
    25: "Giao thức đánh giá",
    26: "Kết quả headline: Bất biến phép xoay ∆ = 0",
    27: "Bảng xếp hạng: so với 8 baselines",
    28: "Hạn chế & Hướng phát triển",
    29: "Hạn chế: nhìn thẳng vào điểm yếu",
    30: "Hướng phát triển",
    31: "Tài liệu tham khảo",
    33: "Backup: APFD — ví dụ tính tay",
    34: "Backup: Multi-trial protocol (vì sao 30 lần?)",
    35: "Backup: AUC vs APFD — vì sao tách bạch?",
    36: "Backup: Resolution probe (bất biến lấy mẫu)",
    37: "Backup: Chi phí ở quy mô lớn",
}
DIVIDERS = {3, 9, 12, 16, 23, 28}
WHITE = RGBColor(0xFA, 0xFA, 0xFA)
TEAL  = RGBColor(0x23, 0x37, 0x3B)


def first_run(sh):
    for p in sh.text_frame.paragraphs:
        for r in p.runs:
            return r
    return None


def run_color(sh):
    r = first_run(sh)
    try:
        if r is not None and r.font.color and r.font.color.type is not None:
            return str(r.font.color.rgb)
    except Exception:
        pass
    return None


def set_text_keep(sh, text, color, bold, size_pt):
    tf = sh.text_frame
    tf.word_wrap = False
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.color.rgb = color
    r.font.bold = bold
    if size_pt:
        r.font.size = Pt(size_pt)


def delete(sh):
    sh._element.getparent().remove(sh._element)


def is_word_char(ch):
    return ch.isalpha() or ch.isdigit() or ch in "%)]"


def starts_word(ch):
    return ch.isalpha() or ch.isdigit() or ch in "([“\""


prs = Presentation(SRC)
title_fixed = body_spaces = 0

for idx, slide in enumerate(prs.slides, 1):
    shapes = [sh for sh in slide.shapes if sh.has_text_frame and sh.text_frame.text.strip()]

    # ---------- 1/2. TITLE MERGE ----------
    if idx in TITLES:
        if idx in DIVIDERS:
            # divider: all text boxes are the (centered, dark) title
            tboxes = sorted(shapes, key=lambda s: (s.top, s.left))
            if tboxes:
                anchor = tboxes[0]
                top0 = anchor.top
                set_text_keep(anchor, TITLES[idx], TEAL, True, 20)
                anchor.left = Inches(0.6); anchor.width = Inches(5.1); anchor.top = top0
                anchor.text_frame.word_wrap = True
                for s in tboxes[1:]:
                    delete(s)
                title_fixed += 1
        else:
            wboxes = [s for s in shapes if run_color(s) == "FAFAFA" and Emu(s.top).inches < 0.6]
            if wboxes:
                wboxes.sort(key=lambda s: (s.top, s.left))
                anchor = wboxes[0]
                sz = None
                r = first_run(anchor)
                try: sz = r.font.size.pt if r and r.font.size else None
                except Exception: sz = None
                set_text_keep(anchor, TITLES[idx], WHITE, True, sz or 16)
                anchor.left = Inches(0.14); anchor.top = Inches(0.10)
                anchor.width = Inches(6.05); anchor.height = Inches(0.34)
                for s in wboxes[1:]:
                    delete(s)
                title_fixed += 1

    # ---------- 3. BODY SPACE FIX ----------
    # re-read remaining boxes
    boxes = []
    for sh in slide.shapes:
        if not sh.has_text_frame: continue
        t = sh.text_frame.text
        if not t.strip(): continue
        boxes.append({
            "sh": sh, "t": t,
            "top": Emu(sh.top).inches, "left": Emu(sh.left).inches,
            "right": Emu(sh.left).inches + Emu(sh.width).inches,
            "h": Emu(sh.height).inches,
        })
    # only NORMAL-height boxes (>=0.105in excludes sub/superscripts ~0.08-0.09)
    norm = [b for b in boxes if b["h"] >= 0.105 and "\n" not in b["t"]]
    norm.sort(key=lambda b: (round(b["top"], 2), b["left"]))
    for i in range(len(norm) - 1):
        a, b = norm[i], norm[i + 1]
        # SAME baseline only (subscripts sit ~0.05 lower -> excluded)
        if abs(a["top"] - b["top"]) > 0.03:
            continue
        gap = b["left"] - a["right"]
        if gap > 0.05:
            continue  # already separated
        if gap < -0.30:
            continue  # heavy overlap -> leave alone (moving boxes breaks layout)
        ta, tb = a["t"], b["t"]
        if not ta or not tb: continue
        if ta.endswith(" ") or tb.startswith(" "): continue
        if ta.endswith("-") or ta.endswith("–"): continue
        if not is_word_char(ta[-1]): continue
        # RIGHT box must start with a LETTER (excludes table numbers / digits)
        if not tb[0].isalpha(): continue
        # spaces-only: never move boxes (moving breaks mixed prose/math lines)
        r = first_run(b["sh"])
        if r is not None:
            r.text = " " + r.text
            body_spaces += 1

prs.save(DST)
print(f"titles merged: {title_fixed} | body spaces added: {body_spaces}")
print("saved", DST)
