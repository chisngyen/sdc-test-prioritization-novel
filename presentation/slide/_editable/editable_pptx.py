#!/usr/bin/env python
"""Build an EDITABLE RoadFury-to-SE2RoadNet.pptx from the .dc.html deck.

Idea: each <section> is rebuilt as
  (1) a background PNG that keeps every figure (SVG, <img>, CSS charts, tables)
      but has the *prose* text made invisible, plus
  (2) one native, editable text box per prose block, placed at the exact DOM
      geometry (tight text rect via Range), with matching font / size / colour.
Fonts (Be Vietnam Pro, JetBrains Mono) are embedded into the .pptx so it renders
identically on any machine.

Subcommands:
  extract                 dump _editable/blocks/sNN.json for every slide (all text blocks)
  build [--slides a,b,c]  build the editable deck (uses _editable/prose.json if present,
                          else every non-figure text block is treated as prose)

prose.json format: { "0": [eid,eid,...], "3": [...], ... }  (eids per slide to extract)
"""
import argparse, base64, glob, json, os, re, sys, time, zipfile, shutil
from io import BytesIO
from PIL import Image
from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

HERE   = os.path.dirname(os.path.abspath(__file__))
SLIDE  = os.path.join(HERE, os.pardir)                       # presentation/slide
GS     = os.path.join(SLIDE, "google-slides")
DECK   = os.path.join(GS, "RoadFury to SE2RoadNet.dc.html")
OUT    = os.path.join(SLIDE, "RoadFury-to-SE2RoadNet.pptx")
BGDIR  = os.path.join(HERE, "bg")
BLKDIR = os.path.join(HERE, "blocks")
FONTDIR= os.path.join(HERE, "fonts")
PROSE  = os.path.join(HERE, "prose.json")
CHROME = r"C:/Program Files/Google/Chrome/Application/chrome.exe"

# 16:9 canvas: 1920x1080 design px -> 13.333in x 7.5in
EMU_PER_PX = 6350.0          # 13.3333in*914400/1920
PT_PER_PX  = 0.5             # 72*13.3333/1920
W = Inches(13.333)
H = Inches(7.5)
SCALE = 2                    # screenshot device scale for crisp backgrounds

FONTS_LINK = ('<link rel="preconnect" href="https://fonts.googleapis.com">'
  '<link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>'
  '<link href="https://fonts.googleapis.com/css2?family=Be+Vietnam+Pro:ital,wght@0,400;0,500;0,600;0,700;0,800;1,400;1,700'
  '&family=JetBrains+Mono:ital,wght@0,400;0,500;0,700;1,400&display=swap" rel="stylesheet">')

TPL = ('<!doctype html><html><head><meta charset="utf-8">' + FONTS_LINK +
  '<style>html,body{margin:0;padding:0;background:#fff}'
  'section{width:1920px;height:1080px;box-sizing:border-box;overflow:hidden;position:relative}</style>'
  '</head><body>%s</body></html>')

# --- DOM walker: collect every "text leaf" block (skips svg + img subtrees) ---
WALKER_JS = r"""
const INLINE=new Set(['SPAN','B','STRONG','I','EM','A','CODE','SUB','SUP','U','SMALL','MARK','TT','KBD','ABBR','FONT','WBR','BDI','BDO']);
function isSvg(el){return el.namespaceURI==='http://www.w3.org/2000/svg'||el.tagName.toLowerCase()==='svg';}
function rgb(s){const m=(s||'').match(/rgba?\(([^)]+)\)/);if(!m)return null;const p=m[1].split(',').map(x=>parseFloat(x));return{r:Math.round(p[0]),g:Math.round(p[1]),b:Math.round(p[2]),a:p.length>3?p[3]:1};}
function hexOf(c){if(!c)return null;return ((1<<24)+(c.r<<16)+(c.g<<8)+c.b).toString(16).slice(1).toUpperCase();}
function directText(el){for(const n of el.childNodes){if(n.nodeType===3&&n.textContent.trim()!=='')return true;}return false;}
function allInlineNonAbs(el){const ch=[...el.children];if(ch.length===0)return false;return ch.every(c=>{if(!INLINE.has(c.tagName))return false;const p=getComputedStyle(c).position;return p!=='absolute'&&p!=='fixed';});}
function runsOf(el){
  const runs=[];
  function emitRun(t,st){if(t==='')return;const c=rgb(st.color);runs.push({text:t,bold:(parseInt(st.fontWeight)||400)>=600,italic:st.fontStyle==='italic',color:hexOf(c),alpha:c?c.a:1,sizePx:parseFloat(st.fontSize),family:st.fontFamily});}
  function rec(node,st){for(const n of node.childNodes){if(n.nodeType===3){emitRun(n.textContent,st);}else if(n.nodeType===1){if(n.tagName==='BR'){runs.push({br:true});}else{rec(n,getComputedStyle(n));}}}}
  rec(el,getComputedStyle(el));return runs;
}
const out=[];let counter=0;
function emit(el){
  const eid=counter++;el.setAttribute('data-eid',String(eid));
  const cs=getComputedStyle(el);const rect=el.getBoundingClientRect();
  let rr=null;try{const rg=document.createRange();rg.selectNodeContents(el);const b=rg.getBoundingClientRect();if(b.width>0||b.height>0)rr={x:b.x,y:b.y,w:b.width,h:b.height};}catch(e){}
  const padL=parseFloat(cs.paddingLeft)||0,padR=parseFloat(cs.paddingRight)||0,padT=parseFloat(cs.paddingTop)||0,padB=parseFloat(cs.paddingBottom)||0;
  const fs=parseFloat(cs.fontSize);let lh=cs.lineHeight,lhr=1.2;
  if(typeof lh==='string'&&lh.endsWith('px'))lhr=parseFloat(lh)/fs;else if(lh==='normal')lhr=1.2;else if(parseFloat(lh))lhr=parseFloat(lh);
  const ls=parseFloat(cs.letterSpacing);
  out.push({eid,tag:el.tagName,text:el.textContent.replace(/\s+/g,' ').trim(),
    rect:{x:rect.x,y:rect.y,w:rect.width,h:rect.height},range:rr,
    content:{x:rect.x+padL,y:rect.y+padT,w:rect.width-padL-padR,h:rect.height-padT-padB},
    align:cs.textAlign,lineHeight:lhr,sizePx:fs,weight:parseInt(cs.fontWeight)||400,style:cs.fontStyle,
    family:cs.fontFamily,color:hexOf(rgb(cs.color)),alpha:(rgb(cs.color)||{a:1}).a,
    transform:cs.textTransform,letterSpacing:(isNaN(ls)?0:ls),runs:runsOf(el)});
}
function walk(el){
  if(el.nodeType!==1)return;if(isSvg(el)||el.tagName==='IMG')return;
  if((directText(el)||allInlineNonAbs(el))&&el.textContent.trim()!==''){emit(el);return;}
  for(const c of el.children)walk(c);
}
walk(document.body);return out;
"""

HIDE_JS = r"""
const ids=arguments[0];
for(const id of ids){
  const el=document.querySelector('[data-eid="'+id+'"]');if(!el)continue;
  const set=n=>{n.style.setProperty('color','transparent','important');n.style.setProperty('-webkit-text-fill-color','transparent','important');n.style.setProperty('text-shadow','none','important');};
  set(el);el.querySelectorAll('*').forEach(set);
}
return ids.length;
"""

# --- collect <a href> rects so figure/button links become clickable overlays ---
LINKS_JS = r"""
return [...document.querySelectorAll('a[href]')].map(a=>{const r=a.getBoundingClientRect();
  return {href:a.href,x:r.x,y:r.y,w:r.width,h:r.height};}).filter(o=>o.w>1&&o.h>1);
"""

def sections():
    deck = open(DECK, encoding="utf-8").read()
    return re.findall(r'<section[\s\S]*?</section>', deck)

def make_driver():
    o = Options()
    o.add_argument("--headless=new"); o.add_argument("--disable-gpu")
    o.add_argument("--hide-scrollbars"); o.add_argument("--force-device-scale-factor=1")
    o.binary_location = CHROME
    d = webdriver.Chrome(options=o)
    return d

def load_section(d, html, idx):
    tmp = os.path.join(GS, "_etmp_%02d.html" % idx)
    open(tmp, "w", encoding="utf-8").write(TPL % html)
    d.get("file:///" + tmp.replace("\\", "/"))
    d.execute_cdp_cmd("Emulation.setDeviceMetricsOverride",
                      {"width": 1920, "height": 1080, "deviceScaleFactor": SCALE, "mobile": False})
    # wait for webfonts so metrics are correct
    try:
        d.execute_async_script("var cb=arguments[0];document.fonts.ready.then(()=>setTimeout(cb,250));")
    except Exception:
        time.sleep(1.0)
    time.sleep(0.25)
    return tmp

def extract_blocks(d):
    return d.execute_script(WALKER_JS)

def shoot_bg(d, hide_eids, outpng):
    if hide_eids:
        d.execute_script(HIDE_JS, [str(e) for e in hide_eids])
    shot = d.execute_cdp_cmd("Page.captureScreenshot",
        {"clip": {"x": 0, "y": 0, "width": 1920, "height": 1080, "scale": SCALE},
         "captureBeyondViewport": True})
    png = base64.b64decode(shot["data"])
    im = Image.open(BytesIO(png))
    im.save(outpng)
    return outpng

# ---------- prose text assembly ----------
def _collapse(t):
    return re.sub(r'\s+', ' ', t)

def paragraphs_of(block):
    """Split runs into paragraphs on <br>; collapse ws; trim paragraph edges."""
    paras, cur = [], []
    for r in block["runs"]:
        if r.get("br"):
            paras.append(cur); cur = []
        else:
            cur.append(dict(r))
    paras.append(cur)
    up = (block.get("transform") == "uppercase")
    lo = (block.get("transform") == "lowercase")
    out = []
    for runs in paras:
        runs = [r for r in runs if r.get("text", "") != ""]
        for r in runs:
            r["text"] = _collapse(r["text"])
            if up: r["text"] = r["text"].upper()
            elif lo: r["text"] = r["text"].lower()
        if runs:
            runs[0]["text"] = runs[0]["text"].lstrip()
            runs[-1]["text"] = runs[-1]["text"].rstrip()
        runs = [r for r in runs if r.get("text", "") != ""]
        out.append(runs)
    # drop fully-empty trailing/leading paragraphs but keep internal blanks
    while out and not out[0]: out.pop(0)
    while out and not out[-1]: out.pop()
    return out

def map_family(fam):
    f = (fam or "").lower()
    if "jetbrains" in f: return "JetBrains Mono"
    if "be vietnam" in f: return "Be Vietnam Pro"
    first = (fam or "Be Vietnam Pro").split(",")[0].strip().strip("'\"")
    return first or "Be Vietnam Pro"

ALIGN = {"left": PP_ALIGN.LEFT, "start": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
         "right": PP_ALIGN.RIGHT, "end": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}

def add_textbox(slide, block):
    paras = paragraphs_of(block)
    if not any(paras): return
    rr = block.get("range") or block["content"]
    x = max(0, block["content"]["x"]); y = max(0, rr["y"])
    w = max(8, block["content"]["w"]); h = max(8, rr["h"])
    # Horizontal slack so a line that exactly fills its box in Chrome does not
    # spuriously wrap (and overlap the block below) when re-rendered. Widening can
    # only *reduce* wrapping -> safe direction. Keep the alignment anchor fixed.
    fs = block.get("sizePx", 16) or 16
    lh = block.get("lineHeight", 1.2) or 1.2
    nlines = max(1, round(rr.get("h", fs) / (fs * lh)))
    # A single-line block is placed at its exact DOM top (vertical_anchor TOP). Applying
    # the DOM line-height as PPT line_spacing on a one-line paragraph adds spurious leading
    # that pushes the glyph DOWN inside its box -> visible downward drift, and overflow off
    # the slide for bottom-anchored blocks (e.g. the title team list). Use 1.0 for one-liners
    # so the glyph sits at the DOM y; keep the real line-height only where it spans >1 line.
    single = nlines <= 1
    slack = 140 if nlines <= 1 else 16
    al = block.get("align", "left")
    if al == "center":
        x -= slack / 2; w += slack
    elif al in ("right", "end"):
        x -= slack; w += slack
    else:
        w += slack
    x = max(0, x); h = h + fs * 0.35
    tb = slide.shapes.add_textbox(Emu(round(x*EMU_PER_PX)), Emu(round(y*EMU_PER_PX)),
                                  Emu(round(w*EMU_PER_PX)), Emu(round(h*EMU_PER_PX)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    blk_color = block.get("color") or "23373B"
    ls = block.get("letterSpacing", 0) or 0
    first = True
    for runs in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = ALIGN.get(block.get("align", "left"), PP_ALIGN.LEFT)
        try: p.line_spacing = 1.0 if single else float(block.get("lineHeight", 1.2))
        except Exception: pass
        p.space_before = Pt(0); p.space_after = Pt(0)
        if not runs:
            continue
        for r in runs:
            run = p.add_run(); run.text = r["text"]
            run.font.size = Pt(round(r.get("sizePx", block["sizePx"]) * PT_PER_PX, 1))
            run.font.bold = bool(r.get("bold"))
            run.font.italic = bool(r.get("italic"))
            run.font.name = map_family(r.get("family", block.get("family")))
            col = r.get("color") if (r.get("alpha", 1) and r.get("color")) else blk_color
            try: run.font.color.rgb = RGBColor.from_string(col)
            except Exception: run.font.color.rgb = RGBColor.from_string("23373B")
            if ls:
                run.font._rPr.set("spc", str(int(round(ls * PT_PER_PX * 100))))

def add_link_rect(slide, lk):
    """Overlay a fully-transparent (but click-catching) hyperlinked rectangle at the
    DOM geometry of an <a href>, so a baked-into-bg figure/button stays clickable."""
    x = max(0.0, lk["x"]); y = max(0.0, lk["y"])
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Emu(round(x*EMU_PER_PX)), Emu(round(y*EMU_PER_PX)),
        Emu(round(lk["w"]*EMU_PER_PX)), Emu(round(lk["h"]*EMU_PER_PX)))
    sp.shadow.inherit = False
    sp.line.fill.background()
    sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    srgb = sp.fill.fore_color._xFill.find(qn("a:srgbClr"))
    if srgb is not None:
        srgb.append(srgb.makeelement(qn("a:alpha"), {"val": "0"}))   # 100% transparent fill
    sp.click_action.hyperlink.address = lk["href"]

# ---------- font embedding ----------
NS = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main",
      "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
      "ct": "http://schemas.openxmlformats.org/package/2006/content-types",
      "rel": "http://schemas.openxmlformats.org/package/2006/relationships"}
REL_FONT = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/font"

FONT_SET = [
    ("Be Vietnam Pro", {"regular": "BeVietnamPro-Regular.ttf", "bold": "BeVietnamPro-Bold.ttf",
                        "italic": "BeVietnamPro-Italic.ttf", "boldItalic": "BeVietnamPro-BoldItalic.ttf"}),
    ("JetBrains Mono", {"regular": "JetBrainsMono-Regular.ttf", "bold": "JetBrainsMono-Bold.ttf",
                        "italic": "JetBrainsMono-Italic.ttf", "boldItalic": "JetBrainsMono-BoldItalic.ttf"}),
]

def embed_fonts(pptx_path):
    tmp = pptx_path + ".tmp"
    zin = zipfile.ZipFile(pptx_path, "r")
    names = set(zin.namelist())
    pres = zin.read("ppt/presentation.xml")
    rels = zin.read("ppt/_rels/presentation.xml.rels")
    cts  = zin.read("[Content_Types].xml")

    # 1) content types: default for fntdata
    ct = etree.fromstring(cts)
    if not any(d.get("Extension") == "fntdata" for d in ct.findall("ct:Default", NS)):
        d = etree.SubElement(ct, "{%s}Default" % NS["ct"])
        d.set("Extension", "fntdata"); d.set("ContentType", "application/x-fontdata")
    cts_new = etree.tostring(ct, xml_declaration=True, encoding="UTF-8", standalone=True)

    # 2) relationships + font parts
    rl = etree.fromstring(rels)
    existing = [r.get("Id") for r in rl.findall("rel:Relationship", NS)]
    n = 0
    while ("rId900%d" % n) in existing: n += 1
    font_parts = []   # (zipname, bytes)
    embedded = []     # (typeface, {slot: rId})
    fidx = 0
    for typeface, slots in FONT_SET:
        idmap = {}
        for slot, fname in slots.items():
            path = os.path.join(FONTDIR, fname)
            if not os.path.exists(path):
                continue
            fidx += 1
            zipname = "ppt/fonts/font%d.fntdata" % fidx
            rid = "rId90%02d" % fidx
            rel = etree.SubElement(rl, "{%s}Relationship" % NS["rel"])
            rel.set("Id", rid); rel.set("Type", REL_FONT)
            rel.set("Target", "fonts/font%d.fntdata" % fidx)
            font_parts.append((zipname, open(path, "rb").read()))
            idmap[slot] = rid
        if idmap:
            embedded.append((typeface, idmap))
    rels_new = etree.tostring(rl, xml_declaration=True, encoding="UTF-8", standalone=True)

    # 3) presentation.xml: flags + embeddedFontLst
    pr = etree.fromstring(pres)
    pr.set("embedTrueTypeFonts", "1"); pr.set("saveSubsetFonts", "0")
    P = "{%s}" % NS["p"]; R = "{%s}" % NS["r"]
    lst = etree.Element(P + "embeddedFontLst")
    for typeface, idmap in embedded:
        ef = etree.SubElement(lst, P + "embeddedFont")
        fo = etree.SubElement(ef, P + "font"); fo.set("typeface", typeface)
        for slot in ("regular", "bold", "italic", "boldItalic"):
            if slot in idmap:
                e = etree.SubElement(ef, P + slot); e.set(R + "id", idmap[slot])
    # insert after p:notesSz (schema order), else after sldSz / sldIdLst, else append
    anchor = None
    for tag in ("p:notesSz", "p:sldSz", "p:sldIdLst"):
        el = pr.find(tag, NS)
        if el is not None:
            anchor = el; break
    if anchor is not None:
        anchor.addnext(lst)
    else:
        pr.append(lst)
    pres_new = etree.tostring(pr, xml_declaration=True, encoding="UTF-8", standalone=True)

    # 4) rewrite zip
    zout = zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED)
    for item in zin.infolist():
        if item.filename == "ppt/presentation.xml": data = pres_new
        elif item.filename == "ppt/_rels/presentation.xml.rels": data = rels_new
        elif item.filename == "[Content_Types].xml": data = cts_new
        else: data = zin.read(item.filename)
        zout.writestr(item, data)
    for zipname, data in font_parts:
        zout.writestr(zipname, data)
    zin.close(); zout.close()
    os.replace(tmp, pptx_path)
    return [t for t, _ in embedded], fidx

# ---------- sources slide (ported) ----------
def add_sources_slide(prs, blank):
    NAVY = RGBColor(0x23,0x37,0x3B); ORANGE = RGBColor(0xEB,0x81,0x1B); GRAY = RGBColor(0x66,0x66,0x66)
    sources = [
        ("Logo HCMUS", "https://www.hcmus.edu.vn"),
        ("Logo ICST 2026", "https://conf.researchr.org/home/icst-2026"),
        ("Video demo (Google Drive)", "https://drive.google.com/file/d/1JC0NY3qfW-if9cM74Zi3d0VaTcA-1el_/view?usp=sharing"),
        ("Video demo (Facebook)", "https://www.facebook.com/reel/1528755028625867"),
        ("Anh duong & histogram do cong", "https://github.com/christianbirchler-org/sensodat"),
        ("Anh chup mo phong (BeamNG.tech)", "https://www.beamng.tech/"),
    ]
    s = prs.slides.add_slide(blank)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(0.9))
    r = tb.text_frame.paragraphs[0].add_run(); r.text = "Nguon hinh anh / Image sources"
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = NAVY; r.font.name = "Be Vietnam Pro"
    body = s.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(11.6), Inches(4.6))
    tf = body.text_frame; tf.word_wrap = True
    for i,(label,url) in enumerate(sources):
        para = tf.paragraphs[0] if i==0 else tf.add_paragraph(); para.space_after = Pt(12)
        r1 = para.add_run(); r1.text = f"{label}:  "; r1.font.size = Pt(16); r1.font.bold = True; r1.font.color.rgb = NAVY; r1.font.name="Be Vietnam Pro"
        r2 = para.add_run(); r2.text = url; r2.font.size = Pt(16); r2.font.color.rgb = ORANGE; r2.hyperlink.address = url; r2.font.name="JetBrains Mono"
    note = s.shapes.add_textbox(Inches(0.9), Inches(6.4), Inches(11.6), Inches(0.9))
    nr = note.text_frame.paragraphs[0].add_run()
    nr.text = ("So do kien truc la hinh TikZ tu ve (fig_arch, fig_arch_se2); cac bieu do APFD/leaderboard/focal "
               "la SVG/CSS tu ve tu so lieu thi nghiem cua nhom -- khong can dan nguon. Anh duong & histogram do cong "
               "ve tu bo du lieu SensoDat [1]; anh chup mo phong tu BeamNG.tech.")
    nr.font.size = Pt(11); nr.font.italic = True; nr.font.color.rgb = GRAY; nr.font.name="Be Vietnam Pro"

# ---------- commands ----------
def cmd_extract(args):
    secs = sections(); d = make_driver()
    try:
        for i, sec in enumerate(secs):
            tmp = load_section(d, sec, i)
            blocks = extract_blocks(d)
            json.dump(blocks, open(os.path.join(BLKDIR, "s%02d.json" % i), "w", encoding="utf-8"),
                      ensure_ascii=False, indent=1)
            os.remove(tmp)
            print("s%02d  %d blocks" % (i, len(blocks)))
    finally:
        d.quit()

def cmd_build(args):
    secs = sections()
    idxs = [int(x) for x in args.slides.split(",")] if args.slides else list(range(len(secs)))
    prose_map = json.load(open(PROSE, encoding="utf-8")) if os.path.exists(PROSE) else None
    out_path = args.out or OUT
    prs = Presentation(); prs.slide_width = W; prs.slide_height = H
    blank = prs.slide_layouts[6]
    d = make_driver()
    try:
        for i in idxs:
            tmp = load_section(d, secs[i], i)
            blocks = extract_blocks(d)
            links = d.execute_script(LINKS_JS)
            bmap = {b["eid"]: b for b in blocks}
            if prose_map is not None:
                prose_eids = [e for e in prose_map.get(str(i), []) if e in bmap]
            else:
                prose_eids = list(bmap.keys())     # default: every text block
            # safety: never lift near-invisible decorative/watermark text out of the image
            prose_eids = [e for e in prose_eids if (bmap[e].get("alpha") or 1) >= 0.12]
            bg = shoot_bg(d, prose_eids, os.path.join(BGDIR, "s-%02d.png" % i))
            os.remove(tmp)
            s = prs.slides.add_slide(blank)
            s.shapes.add_picture(bg, 0, 0, width=W, height=H)
            for e in prose_eids:
                add_textbox(s, bmap[e])
            for lk in links:
                add_link_rect(s, lk)
            print("s%02d  bg+%d boxes+%d links" % (i, len(prose_eids), len(links)))
    finally:
        d.quit()
    if args.slides is None:
        add_sources_slide(prs, blank)
    prs.save(out_path)
    fams, nf = embed_fonts(out_path)
    print("Saved:", os.path.normpath(out_path), "| slides:", len(prs.slides._sldIdLst),
          "| embedded fonts:", fams, "(%d files)" % nf)

if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    sub = ap.add_subparsers(dest="cmd", required=True)
    sub.add_parser("extract")
    b = sub.add_parser("build"); b.add_argument("--slides", default=None); b.add_argument("--out", default=None)
    args = ap.parse_args()
    {"extract": cmd_extract, "build": cmd_build}[args.cmd](args)
