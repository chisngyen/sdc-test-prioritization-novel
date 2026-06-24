#!/usr/bin/env python
"""Build RoadFury-to-SE2RoadNet.pptx from the rendered .dc.html pages (image-faithful),
plus a final clickable 'Sources' slide for every raster image used in the deck.

Render first:  cd ../google-slides && node render-slide.js all "RoadFury to SE2RoadNet.dc.html" "s-"
Then:          python build_pptx.py
"""
import glob, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

HERE   = os.path.dirname(os.path.abspath(__file__))
SLIDES = os.path.join(HERE, os.pardir, "google-slides")     # PNGs land here (s-NN.png)
OUT    = os.path.join(HERE, os.pardir, "RoadFury-to-SE2RoadNet.pptx")

W = Inches(13.333)   # 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

imgs = sorted(glob.glob(os.path.join(SLIDES, "s-*.png")))
if not imgs:
    raise SystemExit("No s-*.png found in %s -- run render-slide.js first." % SLIDES)
for img in imgs:
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(img, 0, 0, width=W, height=H)

# ---- final 'Sources' slide: every raster-image link, clickable (prof requirement) ----
NAVY   = RGBColor(0x23, 0x37, 0x3B)
ORANGE = RGBColor(0xEB, 0x81, 0x1B)
GRAY   = RGBColor(0x66, 0x66, 0x66)

sources = [
    ("Logo HCMUS",                          "https://www.hcmus.edu.vn"),
    ("Logo ICST 2026",                      "https://conf.researchr.org/home/icst-2026"),
    ("Video demo (Google Drive)",           "https://drive.google.com/file/d/1JC0NY3qfW-if9cM74Zi3d0VaTcA-1el_/view?usp=sharing"),
    ("Video demo (Facebook)",               "https://www.facebook.com/reel/1528755028625867"),
    ("Anh duong & histogram do cong",       "https://github.com/christianbirchler-org/sensodat"),
    ("Anh chup mo phong (BeamNG.tech)",     "https://www.beamng.tech/"),
]

s = prs.slides.add_slide(blank)
tb = s.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(0.9))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Nguon hinh anh / Image sources"
r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = NAVY

body = s.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(11.6), Inches(4.6))
tf = body.text_frame; tf.word_wrap = True
for i, (label, url) in enumerate(sources):
    para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    para.space_after = Pt(12)
    r1 = para.add_run(); r1.text = f"{label}:  "
    r1.font.size = Pt(16); r1.font.bold = True; r1.font.color.rgb = NAVY
    r2 = para.add_run(); r2.text = url
    r2.font.size = Pt(16); r2.font.color.rgb = ORANGE
    r2.hyperlink.address = url

note = s.shapes.add_textbox(Inches(0.9), Inches(6.4), Inches(11.6), Inches(0.9))
np = note.text_frame; np.word_wrap = True
nr = np.paragraphs[0].add_run()
nr.text = ("So do kien truc la hinh TikZ tu ve (fig_arch, fig_arch_se2); cac bieu do APFD/leaderboard/focal "
           "la SVG/CSS tu ve tu so lieu thi nghiem cua nhom -- khong can dan nguon. Anh duong & histogram do cong "
           "ve tu bo du lieu SensoDat [1]; anh chup mo phong tu BeamNG.tech. Thumbnail video se duoc thay bang khung hinh video cua nhom.")
nr.font.size = Pt(11); nr.font.italic = True; nr.font.color.rgb = GRAY

prs.save(OUT)
print("Saved:", os.path.normpath(OUT))
print("Slides:", len(prs.slides._sldIdLst))
